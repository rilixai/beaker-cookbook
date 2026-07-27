"""Per-task on-disk staging area for the Harvey LAB agent.

The agent itself does all of its file work inside a Stirrup code-execution
environment (one ``code_exec`` tool, as LAB-AA does), so this module is *not*
an agent toolbelt — it is the host-side staging area either end of that
environment:

* ``documents/`` — the task's input documents, materialized locally and then
  uploaded into the execution environment at session start.
* ``output/`` — where deliverables are pulled back *out* of the environment
  after the agent finishes, so the rubric judge can read them.

Because the agent can now emit real binary deliverables (``.docx`` /
``.xlsx`` / ``.pptx`` built with python-docx / openpyxl / python-pptx inside
the environment), :meth:`TaskWorkspace.collect_deliverables` extracts text
from whatever it finds rather than assuming UTF-8 — LAB-AA grades "the text
extracted from the criterion's declared deliverable files".

The parsers (python-docx / openpyxl / pypdf / python-pptx) are imported
lazily so the module stays importable offline; hermetic tests stage plain
text fixtures and never touch the heavy parsers.
"""

from __future__ import annotations

import email
import logging
import shutil
from collections.abc import Callable, Sequence
from pathlib import Path
from typing import Any


logger = logging.getLogger(__name__)


# Per-case factory: ``(record) -> TaskWorkspace``. Backed by a local
# ``harvey-labs`` checkout; tests inject a fixture-backed one.
TaskSource = Callable[[Any], "TaskWorkspace"]


__all__ = [
    "TaskSource",
    "TaskWorkspace",
    "extract_text",
    "task_source_from_dir",
]


class TaskWorkspace:
    """A materialized ``documents/`` + ``output/`` staging tree for one task."""

    def __init__(self, root: str | Path) -> None:
        self._root = Path(root).resolve()
        self._documents = self._root / "documents"
        self._output = self._root / "output"
        self._documents.mkdir(parents=True, exist_ok=True)
        self._output.mkdir(parents=True, exist_ok=True)

    @property
    def root(self) -> Path:
        return self._root

    @property
    def documents_dir(self) -> Path:
        return self._documents

    @property
    def output_dir(self) -> Path:
        return self._output

    def close(self) -> None:
        """Remove the workspace tree."""
        shutil.rmtree(self._root, ignore_errors=True)

    # ─── path safety ──────────────────────────────────────────────────

    def _resolve_within(self, base: Path, rel_path: str) -> Path:
        candidate = (base / rel_path).resolve()
        if candidate != base and base not in candidate.parents:
            raise ValueError(f"Path {rel_path!r} escapes {base.name}/.")
        return candidate

    # ─── deliverable retrieval (output/) ──────────────────────────────

    def deliverable_path(self, name: str) -> Path:
        """Absolute path a deliverable named ``name`` is pulled back down to."""
        return self._resolve_within(self._output, name)

    def collect_deliverables(self, names: Sequence[str] | None = None) -> dict[str, str]:
        """Extract text from the deliverables pulled back into ``output/``.

        ``names`` restricts collection to the exact filenames the task asked
        for (LAB-AA requires exact-filename matches). Passing ``None`` walks the
        whole ``output/`` tree instead — useful for inspection, not for grading.
        Files that are absent are simply omitted; the caller decides what a
        missing deliverable means for the rubric.
        """
        candidates: list[tuple[str, Path]] = []
        if names is None:
            candidates = [
                (path.relative_to(self._output).as_posix(), path)
                for path in sorted(self._output.rglob("*"))
                if path.is_file()
            ]
        else:
            for name in names:
                try:
                    candidates.append((name, self._resolve_within(self._output, name)))
                except ValueError:
                    logger.warning("Deliverable name %r escapes output/; skipped.", name)
        out: dict[str, str] = {}
        for name, path in candidates:
            if path.is_file():
                out[name] = extract_text(path)
        return out


# ─── lazy document parsers ────────────────────────────────────────────


def extract_text(path: Path | str, *, max_chars: int | None = None) -> str:
    """Extract gradable text from a file, dispatching on its extension.

    Handles the deliverable formats LAB tasks ask for (.docx / .xlsx / .pptx /
    .pdf / .eml) and falls back to a lenient UTF-8 decode for .md and other
    plain text. A parser failure returns a short note rather than raising, so
    one unreadable deliverable never aborts a grading run — the judge simply
    sees that the file could not be read.
    """
    path = Path(path)
    suffix = path.suffix.lower()
    try:
        if suffix == ".docx":
            text = _read_docx(path)
        elif suffix in (".xlsx", ".xlsm"):
            text = _read_xlsx(path)
        elif suffix == ".pptx":
            text = _read_pptx(path)
        elif suffix == ".pdf":
            text = _read_pdf(path)
        elif suffix == ".eml":
            text = _read_eml(path)
        else:
            text = path.read_bytes().decode("utf-8", errors="replace")
    except Exception as exc:  # noqa: BLE001 - a bad file must not abort grading
        logger.warning("Could not extract text from %s: %s", path, exc)
        text = f"Could not extract text from {path.name!r} ({type(exc).__name__}: {exc})."
    return text if max_chars is None else text[:max_chars]


def _read_docx(path: Path) -> str:
    try:
        from docx import Document
    except ImportError as exc:  # pragma: no cover - import guard
        raise ImportError("Text extraction requires the optional `python-docx` dependency for .docx.") from exc
    try:
        doc = Document(str(path))
    except Exception as exc:  # pragma: no cover - defensive
        return f"Could not parse {path.name!r} as .docx ({type(exc).__name__}: {exc})."
    parts: list[str] = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append("\t".join(cells))
    return "\n".join(parts)


def _read_xlsx(path: Path) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:  # pragma: no cover - import guard
        raise ImportError("Text extraction requires the optional `openpyxl` dependency for .xlsx.") from exc
    workbook = load_workbook(filename=str(path), read_only=True, data_only=True)
    try:
        names = list(workbook.sheetnames)
        chunks: list[str] = ["# Sheets: " + " | ".join(names)]
        for name in names:
            chunks.append(f"# Sheet: {name}")
            for row_index, row in enumerate(workbook[name].iter_rows(values_only=True)):
                if row_index >= 200:
                    chunks.append("... (truncated rows)")
                    break
                chunks.append("\t".join("" if c is None else str(c) for c in row[:40]))
        return "\n".join(chunks)
    finally:
        workbook.close()


def _read_pdf(path: Path) -> str:
    head = path.read_bytes()[:5]
    if head[:4] == b"PK\x03\x04":
        return f"{path.name!r} is a ZIP/Office container mislabeled .pdf; read it as .docx/.xlsx instead."
    try:
        from pypdf import PdfReader
    except ImportError as exc:  # pragma: no cover - import guard
        raise ImportError("Text extraction requires the optional `pypdf` dependency for .pdf.") from exc
    try:
        reader = PdfReader(str(path))
    except Exception as exc:  # pragma: no cover - defensive
        return f"Could not parse {path.name!r} as a PDF ({type(exc).__name__}: {exc})."
    parts: list[str] = []
    for page_index, page in enumerate(reader.pages):
        if page_index >= 50:
            parts.append("... (truncated pages)")
            break
        try:
            parts.append(page.extract_text() or "")
        except Exception:  # noqa: BLE001 - defensive per-page
            parts.append("")
    return "\n".join(parts)


def _read_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:  # pragma: no cover - import guard
        raise ImportError("Text extraction requires the optional `python-pptx` dependency for .pptx.") from exc
    presentation = Presentation(str(path))
    parts: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts.append(f"# Slide {index}")
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if isinstance(text, str) and text.strip():
                parts.append(text)
            table = getattr(shape, "table", None) if getattr(shape, "has_table", False) else None
            if table is not None:
                for row in table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        parts.append("\t".join(cells))
    return "\n".join(parts)


def _read_eml(path: Path) -> str:
    msg = email.message_from_bytes(path.read_bytes())
    headers = [f"{h}: {msg.get(h, '')}" for h in ("From", "To", "Cc", "Date", "Subject") if msg.get(h)]
    body_parts: list[str] = []
    if msg.is_multipart():
        for part in msg.walk():
            if part.get_content_type() == "text/plain":
                payload = part.get_payload(decode=True)
                if isinstance(payload, bytes):
                    body_parts.append(payload.decode("utf-8", errors="replace"))
    else:
        payload = msg.get_payload(decode=True)
        if isinstance(payload, bytes):
            body_parts.append(payload.decode("utf-8", errors="replace"))
    return "\n".join(headers) + "\n\n" + "\n".join(body_parts)


# ─── task-source factories ────────────────────────────────────────────


def task_source_from_dir(tasks_root: str | Path) -> TaskSource:
    """Build a ``(record) -> TaskWorkspace`` factory backed by a local tree.

    ``tasks_root`` is the ``tasks/`` directory of a ``harveyai/harvey-labs``
    checkout; a task's ``documents/`` live at ``tasks_root/<task_id>/documents``
    (``task_id`` is the task directory's path relative to ``tasks/``, possibly
    nested under a sub-category). The staged ``documents/`` tree is what gets
    uploaded into the agent's execution environment. Fully offline (used by the
    CLI + tests, which point it at a fixture tree).
    """
    base = Path(tasks_root)

    def _factory(record: Any) -> TaskWorkspace:
        import tempfile

        task_id = str(getattr(record, "task_id", "") or "")
        src_docs = base / task_id / "documents"
        ws = TaskWorkspace(tempfile.mkdtemp(prefix="harvey_lab_"))
        # Clean up the temp tree if the copy fails partway (don't leak on error).
        try:
            if src_docs.is_dir():
                shutil.copytree(src_docs, ws.documents_dir, dirs_exist_ok=True)
        except BaseException:
            ws.close()
            raise
        return ws

    return _factory
